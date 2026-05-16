"""
신용카드 이상거래 탐지 프로젝트 PPT 생성 스크립트
- 데이터 특성 및 사기 패턴 슬라이드 포함
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 상수 ────────────────────────────────────────────────
NAVY    = RGBColor(0x1F, 0x35, 0x64)   # 진남색 (헤더/강조)
BLUE    = RGBColor(0x2E, 0x75, 0xB6)   # 파란색 (소제목)
RED     = RGBColor(0xC0, 0x00, 0x00)   # 빨간색 (사기/경고 강조)
ORANGE  = RGBColor(0xED, 0x7D, 0x31)   # 주황색 (포인트)
GRAY    = RGBColor(0x59, 0x59, 0x59)   # 회색 (본문)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)  # 연파랑 (배경 박스)
LIGHT_RED  = RGBColor(0xFC, 0xE4, 0xD6)  # 연주황 (사기 강조 박스)
GREEN   = RGBColor(0x37, 0x86, 0x44)   # 초록 (정상/좋음)


# ── 헬퍼 함수 ────────────────────────────────────────────────

def add_slide(prs, layout_index=6):
    """빈 슬라이드 추가"""
    layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(layout)


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=GRAY,
                align=PP_ALIGN.LEFT, bg_color=None, wrap=True):
    """텍스트박스 추가"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = '맑은 고딕'

    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """사각형 도형 추가"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_slide_header(slide, title, subtitle=None):
    """슬라이드 상단 헤더 바 추가"""
    # 진남색 헤더 바
    add_rect(slide, 0, 0, 10, 1.1, NAVY)
    # 제목
    add_textbox(slide, title, 0.3, 0.1, 9.4, 0.8,
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # 부제목
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.75, 9.4, 0.4,
                    font_size=13, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)


def add_bullet_box(slide, items, left, top, width, height,
                   title=None, title_color=BLUE, bg=LIGHT_BLUE, font_size=14):
    """항목 목록 박스 추가"""
    add_rect(slide, left, top, width, height, bg)
    y = top + 0.1
    if title:
        add_textbox(slide, title, left + 0.15, y, width - 0.2, 0.4,
                    font_size=15, bold=True, color=title_color)
        y += 0.38
    for item in items:
        add_textbox(slide, item, left + 0.15, y, width - 0.3, 0.42,
                    font_size=font_size, color=GRAY)
        y += 0.38


# ────────────────────────────────────────────────────────────
# 슬라이드 제작 시작
# ────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)


# ══════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ══════════════════════════════════════════════════════════════
slide = add_slide(prs)

# 배경 전체를 진남색으로
add_rect(slide, 0, 0, 10, 7.5, NAVY)

# 장식 선
add_rect(slide, 0.5, 2.5, 9, 0.05, BLUE)
add_rect(slide, 0.5, 5.1, 9, 0.05, BLUE)

# 메인 타이틀
add_textbox(slide, '💳 신용카드 이상거래 탐지', 0.5, 2.7, 9, 0.9,
            font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 'Credit Card Fraud Detection', 0.5, 3.55, 9, 0.6,
            font_size=20, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# 부제목
add_textbox(slide, '머신러닝 기반 불균형 이진 분류 프로젝트', 0.5, 4.25, 9, 0.5,
            font_size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# 하단 정보
add_textbox(slide, '284,807건 거래   |   사기 비율 0.17%   |   9개 모델 비교', 0.5, 5.3, 9, 0.4,
            font_size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_textbox(slide, 'ksh1004', 0.5, 6.7, 9, 0.4,
            font_size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# 슬라이드 2: 프로젝트 개요
# ══════════════════════════════════════════════════════════════
slide = add_slide(prs)
add_slide_header(slide, '프로젝트 개요', '왜 이 문제가 어려운가?')

# 핵심 문제 박스
add_rect(slide, 0.4, 1.25, 9.2, 1.5, LIGHT_RED)
add_textbox(slide, '핵심 문제', 0.6, 1.3, 9, 0.38,
            font_size=14, bold=True, color=RED)
add_textbox(slide,
    '"전부 정상" 이라고만 예측해도 정확도 99.8%\n'
    '→ Accuracy는 의미 없음. F1-Score · AUC-ROC · Recall로 평가해야 합니다.',
    0.6, 1.65, 8.8, 0.9, font_size=15, color=NAVY)

# 3개 목표 박스
for i, (icon, title, desc) in enumerate([
    ('🎯', '목표', '사기 거래를 높은 Recall로 탐지'),
    ('📊', '데이터', 'Kaggle Credit Card Fraud (284,807건)'),
    ('🤖', '방법론', '9개 모델 비교 + SMOTE 불균형 처리'),
]):
    x = 0.4 + i * 3.1
    add_rect(slide, x, 2.95, 2.9, 1.6, LIGHT_BLUE)
    add_textbox(slide, icon + ' ' + title, x + 0.15, 3.0, 2.6, 0.4,
                font_size=15, bold=True, color=NAVY)
    add_textbox(slide, desc, x + 0.15, 3.42, 2.6, 1.0, font_size=13, color=GRAY)

# 평가지표 설명
add_rect(slide, 0.4, 4.75, 9.2, 2.35, RGBColor(0xF2, 0xF2, 0xF2))
add_textbox(slide, '📐 평가 지표', 0.6, 4.82, 9, 0.38,
            font_size=14, bold=True, color=NAVY)

metrics = [
    ('Recall(탐지율)', '실제 사기 중 모델이 탐지한 비율 → 놓치면 금전 피해, 가장 중요'),
    ('Precision(정밀도)', '사기 예측 중 실제 사기 비율 → 오탐 시 고객 불편'),
    ('F1-Score', 'Precision + Recall 조화 평균 → 불균형 데이터의 종합 지표'),
    ('AUC-ROC', '임계값 무관 전체 판별력 → 1에 가까울수록 좋음'),
]
for j, (name, desc) in enumerate(metrics):
    y = 5.25 + j * 0.35
    add_textbox(slide, f'• {name}:', 0.6, y, 2.3, 0.34,
                font_size=12, bold=True, color=BLUE)
    add_textbox(slide, desc, 2.9, y, 6.5, 0.34, font_size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════
# 슬라이드 3: 데이터 소개
# ══════════════════════════════════════════════════════════════
slide = add_slide(prs)
add_slide_header(slide, '데이터 소개', 'Kaggle - Credit Card Fraud Detection Dataset')

# 왼쪽: 데이터 기본 정보
add_rect(slide, 0.4, 1.25, 4.3, 5.8, LIGHT_BLUE)
add_textbox(slide, '📋 데이터 기본 정보', 0.6, 1.32, 4.0, 0.4,
            font_size=14, bold=True, color=NAVY)

info_items = [
    ('출처', 'Kaggle (ULB 머신러닝 그룹)'),
    ('총 거래 건수', '284,807건'),
    ('사기 건수', '492건 (0.17%)'),
    ('정상 건수', '284,315건 (99.83%)'),
    ('특성 수', '30개 (V1~V28, Amount, Time)'),
    ('타겟', 'Class (0=정상, 1=사기)'),
    ('수집 기간', '2013년 9월 (이틀간)'),
]
for k, (key, val) in enumerate(info_items):
    y = 1.82 + k * 0.62
    add_textbox(slide, key, 0.6, y, 1.6, 0.38,
                font_size=12, bold=True, color=BLUE)
    add_textbox(slide, val, 2.2, y, 2.3, 0.38, font_size=12, color=GRAY)

# 오른쪽: 컬럼 설명
add_rect(slide, 4.9, 1.25, 4.7, 5.8, RGBColor(0xF2, 0xF2, 0xF2))
add_textbox(slide, '🔍 컬럼 설명', 5.1, 1.32, 4.3, 0.4,
            font_size=14, bold=True, color=NAVY)

cols = [
    ('Time', '첫 거래로부터 경과 시간(초)'),
    ('V1 ~ V28', 'PCA로 익명화된 거래 특성\n원본 컬럼명은 개인정보 보호로 비공개'),
    ('Amount', '거래 금액 (달러)'),
    ('Class', '0 = 정상 거래\n1 = 사기 거래 (타겟)'),
]
y_col = 1.82
for col_name, col_desc in cols:
    add_textbox(slide, col_name, 5.1, y_col, 1.3, 0.35,
                font_size=12, bold=True, color=RED if col_name == 'Class' else BLUE)
    add_textbox(slide, col_desc, 6.4, y_col, 2.9, 0.55, font_size=11, color=GRAY)
    y_col += 0.72 if '\n' in col_desc else 0.55

add_textbox(slide, '⚠️ V1~V28은 이미 PCA 정규화 완료\nAmount, Time만 추가 정규화 필요',
            5.1, 5.5, 4.3, 0.85, font_size=12, color=ORANGE)


# ══════════════════════════════════════════════════════════════
# 슬라이드 4: 데이터 특성 분석
# ══════════════════════════════════════════════════════════════
slide = add_slide(prs)
add_slide_header(slide, '데이터 특성 분석', '극심한 클래스 불균형 — 사기 거래는 0.17%')

# 불균형 시각화 (텍스트 기반 바 차트)
add_rect(slide, 0.4, 1.25, 9.2, 1.8, LIGHT_BLUE)
add_textbox(slide, '📊 클래스 분포 (전체 284,807건)', 0.6, 1.3, 9, 0.38,
            font_size=14, bold=True, color=NAVY)

# 정상 거래 바
add_rect(slide, 0.6, 1.82, 8.5, 0.4, GREEN)
add_textbox(slide, '정상 거래  284,315건  (99.83%)', 0.75, 1.87, 8.2, 0.32,
            font_size=13, bold=True, color=WHITE)

# 사기 거래 바 (매우 얇게)
add_rect(slide, 0.6, 2.32, 0.015, 0.4, RED)
add_rect(slide, 0.615, 2.32, 8.485, 0.4, RGBColor(0xE0, 0xE0, 0xE0))
add_textbox(slide, '사기 거래  492건  (0.17%) ◀ 이것이 탐지 대상', 0.75, 2.37, 8.2, 0.32,
            font_size=13, bold=True, color=RED)

# 3개 특성 카드
card_data = [
    ('💰 거래 금액 패턴', RED, LIGHT_RED, [
        '• 정상 거래 평균: $88',
        '• 사기 거래 평균: $122',
        '• 사기는 소액(1~100달러) 비중 높음',
        '• 소액 테스트 후 고액 결제 패턴 존재',
        '• 정상 거래는 넓은 금액대에 분산',
    ]),
    ('⏰ 시간대 패턴', BLUE, LIGHT_BLUE, [
        '• 새벽 시간대(0~6시) 사기 집중',
        '• 정상 거래는 낮 시간대 집중',
        '• 카드 도용 후 빠른 연속 결제',
        '• Time 컬럼: 0~172,792초(약 48시간)',
        '• 새벽 거래는 이상 탐지 신호',
    ]),
    ('📐 피처 분포 패턴', ORANGE, RGBColor(0xFF, 0xF2, 0xCC), [
        '• V14, V17, V12가 사기 탐지 핵심 피처',
        '• 사기 거래는 V14값이 매우 낮음(-10 이하)',
        '• 정상 거래의 V 피처는 0 근방에 집중',
        '• 사기는 극단값(양/음 모두)에 위치',
        '• PCA 변환으로 해석은 제한적',
    ]),
]
for i, (title, title_color, bg_color, items) in enumerate(card_data):
    x = 0.4 + i * 3.1
    add_rect(slide, x, 3.0, 2.95, 4.1, bg_color)
    add_textbox(slide, title, x + 0.12, 3.07, 2.7, 0.4,
                font_size=13, bold=True, color=title_color)
    for j, item in enumerate(items):
        add_textbox(slide, item, x + 0.12, 3.55 + j * 0.64, 2.7, 0.58,
                    font_size=11, color=GRAY)


# ══════════════════════════════════════════════════════════════
# 슬라이드 5: 사기 패턴 심층 분석
# ══════════════════════════════════════════════════════════════
slide = add_slide(prs)
add_slide_header(slide, '사기 패턴 심층 분석', '왜 사기 거래는 극소수이며, 어떤 특징을 가지는가?')

# 왜 극소수인가
add_rect(slide, 0.4, 1.25, 9.2, 1.1, LIGHT_RED)
add_textbox(slide, '❓ 왜 사기는 0.17%밖에 없을까?', 0.6, 1.3, 9, 0.38,
            font_size=14, bold=True, color=RED)
add_textbox(slide,
    '카드사의 실시간 감지 시스템으로 대부분 차단되어 데이터에 남는 건수가 적고, '
    '사기범도 들키지 않기 위해 정상처럼 보이는 소규모 거래를 선택합니다.',
    0.6, 1.68, 8.8, 0.55, font_size=13, color=NAVY)

# 4가지 사기 패턴
pattern_data = [
    ('💸', '소액 테스트 패턴', LIGHT_RED, RED,
     '카드 도용 직후 1~10달러 소액으로\n사용 가능 여부를 먼저 테스트.\n성공하면 고액 결제로 이어짐.'),
    ('🌙', '새벽 시간대 집중', LIGHT_BLUE, BLUE,
     '자동화 도구를 이용한 사기가\n새벽 0~6시에 집중.\n모니터링 공백을 노린 패턴.'),
    ('⚡', '짧은 시간 내 반복', RGBColor(0xFF, 0xF2, 0xCC), ORANGE,
     '수 분~수십 분 사이에 같은 카드로\n여러 번 결제. 정상 소비 행동과\n다른 이상 빈도 패턴.'),
    ('📉', 'V14 극단값 분포', RGBColor(0xE2, 0xEF, 0xDA), GREEN,
     '모델이 찾아낸 핵심 패턴.\n사기 거래의 V14 값은 -10 이하로\n정상 범위(-3~3)를 크게 벗어남.'),
]

for i, (icon, title, bg, fg, desc) in enumerate(pattern_data):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 4.65
    y = 2.55 + row * 2.3
    add_rect(slide, x, y, 4.4, 2.1, bg)
    add_textbox(slide, icon + '  ' + title, x + 0.15, y + 0.08, 4.1, 0.42,
                font_size=14, bold=True, color=fg)
    add_textbox(slide, desc, x + 0.15, y + 0.55, 4.1, 1.4, font_size=12, color=GRAY)

# 하단 핵심 메시지
add_rect(slide, 0.4, 7.05, 9.2, 0.35, NAVY)
add_textbox(slide, '→ 모델은 이 미묘한 패턴 차이를 통계적으로 학습해 사기를 구별합니다.',
            0.6, 7.07, 9, 0.3, font_size=12, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════
# 슬라이드 6: 전처리 과정
# ══════════════════════════════════════════════════════════════
slide = add_slide(prs)
add_slide_header(slide, '전처리 파이프라인', '데이터를 모델이 학습하기 좋은 형태로 변환')

steps = [
    ('1', '정규화 (StandardScaler)', BLUE, LIGHT_BLUE,
     'Amount($0~$25,000)와 Time(초) 스케일이 V1~V28과 달라\n모델이 큰 값에 편향될 수 있음\n→ 각 컬럼에 독립된 Scaler 적용 (평균=0, 표준편차=1)',
     '⚠️ 컬럼마다 독립 Scaler — 하나의 Scaler로 순서대로 처리하면\n   두 번째 fit()이 첫 번째 기준을 덮어씀 (버그 발생 지점)'),
    ('2', '훈련/테스트 분리 (8:2)', BLUE, LIGHT_BLUE,
     '전체 데이터의 80%는 학습용, 20%는 성능 측정용으로 분리\n→ 테스트 데이터는 학습에 절대 사용하지 않음 (현실 시뮬레이션)',
     'stratify=y 옵션: 분리 후에도 사기 비율(0.17%) 동일 유지\n   없으면 사기가 한쪽에 몰릴 수 있음'),
    ('3', 'SMOTE (불균형 해소)', RED, LIGHT_RED,
     '훈련 데이터: 정상 227,451건 vs 사기 394건 → 약 578:1 불균형\n→ SMOTE로 사기 데이터를 227,451건으로 증강 (1:1 균형)',
     '⚠️ 테스트 데이터에는 절대 적용 금지\n   현실에 없는 가상 데이터로 평가하면 성능이 왜곡됨'),
]

for i, (num, title, color, bg, desc, warn) in enumerate(steps):
    y = 1.25 + i * 1.9
    add_rect(slide, 0.4, y, 9.2, 1.78, bg)
    # 번호 원
    add_rect(slide, 0.5, y + 0.15, 0.55, 0.55, color)
    add_textbox(slide, num, 0.5, y + 0.1, 0.55, 0.55,
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 제목
    add_textbox(slide, title, 1.2, y + 0.12, 8, 0.4,
                font_size=15, bold=True, color=color)
    # 설명
    add_textbox(slide, desc, 1.2, y + 0.52, 7.8, 0.65, font_size=11, color=GRAY)
    # 경고/포인트
    add_textbox(slide, warn, 1.2, y + 1.15, 7.8, 0.52, font_size=10, color=ORANGE)


# ══════════════════════════════════════════════════════════════
# 슬라이드 7: 기초 모델 설명
# ══════════════════════════════════════════════════════════════
slide = add_slide(prs)
add_slide_header(slide, '사용 모델 ① — 기초 모델', '5가지 고전 머신러닝 분류 알고리즘')

basic_models = [
    ('로지스틱 회귀', '각 피처에 가중치를 곱해 합산 → 시그모이드 함수로 0~1 확률 변환 → 0.5 기준으로 정상/사기 판정', '빠르고 해석이 쉬움. 선형 경계만 학습 가능'),
    ('결정 트리', '"V14 < -5?" 같은 조건 질문을 반복해 트리를 쌓아 분류. max_depth=10으로 과적합 제한', '직관적. 단독 사용 시 과적합 위험 → 앙상블로 발전'),
    ('KNN', '새 거래와 가장 비슷한 5개 이웃을 찾아 다수결로 정상/사기 판정. 이번 실험 F1 1위(0.859)', '직관적이지만 데이터가 클수록 느려짐 (O(n))'),
    ('나이브 베이즈', '각 피처가 독립이라 가정 + 조건부 확률로 클래스 판정. 매우 빠름', '독립 가정이 맞지 않는 데이터에선 성능 제한'),
    ('SVM', '두 클래스 사이 여백(margin)을 최대화하는 경계 탐색. RBF 커널로 비선형 처리', '고차원에 강하나 O(n²) 복잡도 → 1만 건 샘플로 학습'),
]

for i, (name, desc, note) in enumerate(basic_models):
    y = 1.22 + i * 1.18
    bg = LIGHT_RED if name == 'KNN' else LIGHT_BLUE
    add_rect(slide, 0.4, y, 9.2, 1.1, bg)
    add_textbox(slide, name, 0.6, y + 0.06, 2.2, 0.38,
                font_size=13, bold=True, color=RED if name == 'KNN' else BLUE)
    add_textbox(slide, desc, 2.85, y + 0.06, 6.6, 0.55, font_size=11, color=GRAY)
    add_textbox(slide, '→ ' + note, 2.85, y + 0.62, 6.6, 0.38, font_size=10, color=ORANGE)

add_textbox(slide, '★ KNN이 F1-Score 1위 — 이웃 기반 유사도가 이 데이터의 사기 패턴과 잘 맞음',
            0.4, 7.1, 9.2, 0.32, font_size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# 슬라이드 8: 앙상블 + 딥러닝 모델
# ══════════════════════════════════════════════════════════════
slide = add_slide(prs)
add_slide_header(slide, '사용 모델 ② — 앙상블 & 딥러닝', '여러 모델의 결합과 신경망 구조')

# 앙상블 섹션
add_rect(slide, 0.4, 1.25, 9.2, 3.5, LIGHT_BLUE)
add_textbox(slide, '🌲 앙상블 모델 — 여러 모델을 결합해 더 강한 예측',
            0.6, 1.3, 9, 0.4, font_size=14, bold=True, color=NAVY)

ensemble_data = [
    ('랜덤 포레스트', '결정 트리 100개를 서로 다른 데이터·피처 조합으로 학습 → 다수결\n단일 트리보다 안정적, 과적합에 강함', 'F1: 0.821  AUC: 0.969'),
    ('XGBoost',      '이전 트리의 오류에 집중해 다음 트리를 학습 (Gradient Boosting)\nscale_pos_weight로 불균형 자동 보정 (정상/사기 비율 = 578)', 'F1: 0.602  AUC: 0.978'),
    ('LightGBM',     'XGBoost와 같은 부스팅이지만 리프 단위 성장으로 더 빠름\n금융권 실무에서 속도·성능 균형으로 가장 많이 사용', 'F1: 0.640  AUC: 0.969'),
]

for i, (name, desc, score) in enumerate(ensemble_data):
    y = 1.78 + i * 0.95
    add_textbox(slide, name, 0.6, y, 1.9, 0.38, font_size=13, bold=True, color=BLUE)
    add_textbox(slide, desc, 2.55, y, 5.6, 0.75, font_size=11, color=GRAY)
    add_textbox(slide, score, 8.2, y, 1.3, 0.38, font_size=11, bold=True, color=GREEN)

# 딥러닝 섹션
add_rect(slide, 0.4, 4.9, 9.2, 2.3, RGBColor(0xF2, 0xF2, 0xF2))
add_textbox(slide, '🧠 딥러닝 모델 — 뇌 뉴런 구조를 모방한 신경망',
            0.6, 4.95, 9, 0.4, font_size=14, bold=True, color=NAVY)

# 신경망 구조 시각화 (텍스트)
add_textbox(slide, '신경망(MLP) 구조:', 0.6, 5.45, 2.5, 0.38,
            font_size=13, bold=True, color=BLUE)
add_textbox(slide, '입력층\n(30개 피처)', 3.0, 5.35, 1.5, 0.7,
            font_size=11, color=GRAY, align=PP_ALIGN.CENTER, bg_color=LIGHT_BLUE)
add_textbox(slide, '→', 4.55, 5.55, 0.3, 0.35, font_size=16, bold=True, color=NAVY)
add_textbox(slide, '은닉층 1\n(64 뉴런)', 4.9, 5.35, 1.3, 0.7,
            font_size=11, color=GRAY, align=PP_ALIGN.CENTER, bg_color=LIGHT_BLUE)
add_textbox(slide, '→', 6.25, 5.55, 0.3, 0.35, font_size=16, bold=True, color=NAVY)
add_textbox(slide, '은닉층 2\n(32 뉴런)', 6.6, 5.35, 1.3, 0.7,
            font_size=11, color=GRAY, align=PP_ALIGN.CENTER, bg_color=LIGHT_BLUE)
add_textbox(slide, '→', 7.95, 5.55, 0.3, 0.35, font_size=16, bold=True, color=NAVY)
add_textbox(slide, '출력층\n(정상/사기)', 8.3, 5.35, 1.2, 0.7,
            font_size=11, color=WHITE, align=PP_ALIGN.CENTER, bg_color=NAVY)

add_textbox(slide, 'early_stopping=True: 검증 성능이 개선되지 않으면 조기 종료 → 과적합 방지  |  F1: 0.760  AUC: 0.970',
            0.6, 6.22, 9, 0.38, font_size=11, color=ORANGE)


# ══════════════════════════════════════════════════════════════
# 슬라이드 9: 모델별 학습 전략
# ══════════════════════════════════════════════════════════════
slide = add_slide(prs)
add_slide_header(slide, '모델별 학습 전략', '모든 모델에 동일한 데이터를 쓰지 않는 이유')

# 2가지 전략 박스
add_rect(slide, 0.4, 1.25, 4.45, 5.8, LIGHT_BLUE)
add_textbox(slide, '📗 SMOTE 데이터 사용', 0.6, 1.32, 4.1, 0.4,
            font_size=14, bold=True, color=BLUE)
add_textbox(slide, '로지스틱 회귀\n결정 트리\n랜덤 포레스트\nXGBoost\nLightGBM\n신경망(MLP)',
            0.75, 1.82, 3.8, 2.2, font_size=14, color=NAVY)
add_textbox(slide,
    '이유:\n'
    'SMOTE로 증강된 45만 건 데이터로\n'
    '학습 → 사기 패턴을 더 많이\n'
    '학습할 수 있음\n\n'
    '이 모델들은 대용량 데이터에서도\n충분히 빠르게 학습 가능',
    0.75, 3.95, 3.9, 2.8, font_size=12, color=GRAY)

add_rect(slide, 5.15, 1.25, 4.45, 5.8, LIGHT_RED)
add_textbox(slide, '📕 원본 데이터 사용', 5.35, 1.32, 4.1, 0.4,
            font_size=14, bold=True, color=RED)
add_textbox(slide, 'KNN\n나이브 베이즈\nSVM', 5.5, 1.82, 3.8, 1.3,
            font_size=14, color=NAVY)
add_textbox(slide,
    '이유:\n'
    'KNN: 45만 건 전체를 매 예측마다\n'
    '참조 → 매우 느림 (O(n))\n\n'
    'SVM: O(n²) 복잡도로 23만 건\n'
    '학습 시 수십 분 소요\n'
    '→ 약 1만 건으로 샘플링\n\n'
    '나이브 베이즈: 대용량 증강 데이터가\n'
    '독립 가정을 더 크게 위반',
    5.5, 2.95, 3.9, 3.8, font_size=12, color=GRAY)

add_rect(slide, 0.4, 7.1, 9.2, 0.32, NAVY)
add_textbox(slide, '공통: RANDOM_STATE = 111 고정 → 어떤 환경에서 실행해도 동일한 결과 재현 보장',
            0.6, 7.12, 9, 0.28, font_size=11, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════
# 슬라이드 10: 성능 결과
# ══════════════════════════════════════════════════════════════
slide = add_slide(prs)
add_slide_header(slide, '성능 결과', 'F1-Score 기준 순위 — 불균형 데이터의 핵심 지표')

# 결과 표
headers = ['순위', '모델', 'F1-Score', 'AUC-ROC', 'Recall', 'Precision']
col_w   = [0.55, 1.85, 1.2, 1.2, 1.2, 1.2]
col_x   = [0.3]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

# 헤더 행
y_table = 1.28
for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    add_rect(slide, x, y_table, w - 0.04, 0.4, NAVY)
    add_textbox(slide, h, x + 0.04, y_table + 0.04, w - 0.08, 0.32,
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 데이터 행
rows = [
    ('🥇', 'KNN',          '0.859', '0.944', '0.806', '0.919'),
    ('🥈', '랜덤 포레스트',  '0.821', '0.969', '0.816', '0.825'),
    ('🥉', '신경망(MLP)',   '0.760', '0.970', '0.806', '0.718'),
    ('4',  'LightGBM',     '0.640', '0.969', '0.888', '0.500'),
    ('5',  'XGBoost',      '0.602', '0.978', '0.857', '0.464'),
    ('6',  'SVM',          '0.214', '0.977', '0.898', '0.122'),
    ('7',  '결정 트리',     '0.144', '0.895', '0.806', '0.079'),
    ('8',  '나이브 베이즈', '0.110', '0.963', '0.847', '0.059'),
    ('9',  '로지스틱 회귀', '0.109', '0.970', '0.918', '0.058'),
]

for i, row in enumerate(rows):
    y_row = y_table + 0.44 + i * 0.52
    bg = RGBColor(0xFF, 0xF0, 0xF0) if i == 0 else (
         LIGHT_BLUE if i % 2 == 1 else WHITE)
    for j, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        add_rect(slide, x, y_row, w - 0.04, 0.48, bg)
        c = RED if (i == 0 and j in [0, 1, 2]) else GRAY
        add_textbox(slide, cell, x + 0.04, y_row + 0.07, w - 0.08, 0.35,
                    font_size=11, bold=(i == 0), color=c, align=PP_ALIGN.CENTER)

# 주석
add_textbox(slide,
    '* SVM: 약 1만 건 샘플 학습으로 성능 제약 (전체 학습 시 향상 예상)\n'
    '* XGBoost/LightGBM: 하이퍼파라미터 튜닝 시 F1 0.85+ 달성 가능',
    0.3, 6.93, 9.4, 0.5, font_size=10, color=ORANGE)


# ══════════════════════════════════════════════════════════════
# 슬라이드 11: 결론 및 인사이트
# ══════════════════════════════════════════════════════════════
slide = add_slide(prs)
add_slide_header(slide, '결론 및 핵심 인사이트', '이 프로젝트에서 배운 것들')

insights = [
    (RED,   '💡',
     'Accuracy는 불균형 데이터에서 의미 없다',
     '"전부 정상" 예측도 정확도 99.8%. 반드시 F1-Score · Recall · AUC-ROC로 평가해야 합니다.'),
    (BLUE,  '💡',
     'SMOTE는 강력하지만 올바른 사용법이 중요하다',
     '테스트 데이터에 적용하면 현실과 다른 왜곡된 성능이 나옵니다. 훈련 데이터에만 적용해야 합니다.'),
    (GREEN, '💡',
     '사기 탐지에서는 Recall이 Precision보다 우선이다',
     '사기를 놓치면(FN) 금전 손실. 오탐(FP)은 고객 불편 수준. 비용이 다릅니다.'),
    (ORANGE,'💡',
     '앙상블이 항상 이기지는 않는다',
     '하이퍼파라미터 튜닝 없이는 KNN(기초 모델)이 RF/XGB/LGBM을 F1 기준으로 이겼습니다.'),
    (NAVY,  '💡',
     '모든 모델에 같은 전략을 쓸 수 없다',
     'SVM은 O(n²) 속도 문제로 샘플링 필요. KNN/나이브 베이즈는 SMOTE 데이터가 역효과. 모델별 맞춤 전략이 필요합니다.'),
]

for i, (color, icon, title, desc) in enumerate(insights):
    y = 1.25 + i * 1.2
    bg = RGBColor(0xF8, 0xF8, 0xF8)
    add_rect(slide, 0.4, y, 9.2, 1.1, bg)
    add_rect(slide, 0.4, y, 0.12, 1.1, color)  # 왼쪽 색상 바
    add_textbox(slide, icon + '  ' + title, 0.65, y + 0.06, 8.7, 0.38,
                font_size=13, bold=True, color=color)
    add_textbox(slide, desc, 0.65, y + 0.5, 8.7, 0.5, font_size=12, color=GRAY)

add_rect(slide, 0.4, 7.15, 9.2, 0.28, NAVY)
add_textbox(slide, '데이터의 특성을 이해하고, 모델에 맞는 전략을 선택하는 것이 핵심입니다.',
            0.6, 7.16, 9, 0.25, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# 저장
# ══════════════════════════════════════════════════════════════
output_path = r'c:\Users\Kim SoonHa\Desktop\Code\portfolio_project\fraud_detection\fraud_detection_ppt.pptx'
prs.save(output_path)
print(f'PPT 저장 완료: {output_path}')
print(f'총 슬라이드 수: {len(prs.slides)}장')
